"""
Document loader - downloads and extracts text from documents
"""

import os
import logging
import asyncio
from typing import Optional
from concurrent.futures import ThreadPoolExecutor

try:
    import fitz  # PyMuPDF - much faster than PyPDF2
    PYMUPDF_AVAILABLE = True
except ImportError:
    PYMUPDF_AVAILABLE = False
    from PyPDF2 import PdfReader

from docx import Document
from pptx import Presentation
from app.services.storage_client import storage_client

logger = logging.getLogger(__name__)

# Thread pool for CPU-intensive text extraction
_extraction_executor = ThreadPoolExecutor(max_workers=4)

class DocumentLoader:
    """Load and extract text from documents"""
    
    async def download_and_extract_text(self, file_path: str, document_id: str) -> Optional[str]:
        """
        Download document from storage and extract text
        
        Args:
            file_path: Path in storage bucket (from documents.jsonl_file_path)
            document_id: Document UUID (for temp file naming)
            
        Returns:
            Extracted text content or None if failed
        """
        # Determine file extension
        _, ext = os.path.splitext(file_path)
        ext = ext.lower()
        
        # Create temp directory if not exists
        temp_dir = "/tmp/jsonl_generator"
        os.makedirs(temp_dir, exist_ok=True)
        
        # Download to temp file
        temp_file_path = os.path.join(temp_dir, f"{document_id}{ext}")
        
        try:
            # Download from storage
            success = await storage_client.download_document(file_path, temp_file_path)
            if not success:
                logger.error(f"Failed to download document {file_path}")
                return None
            
            # Extract text based on file type (run in thread pool to avoid blocking)
            file_size_mb = os.path.getsize(temp_file_path) / (1024 * 1024)
            logger.info(f"Extracting text from {ext} file ({file_size_mb:.2f} MB)...")
            
            try:
                if ext == '.pdf':
                    # Run PDF extraction in thread pool with 90 second timeout
                    # Large PDFs (400+ pages) can take 30-60 seconds
                    text = await asyncio.wait_for(
                        asyncio.get_event_loop().run_in_executor(
                            _extraction_executor,
                            self._extract_from_pdf,
                            temp_file_path
                        ),
                        timeout=90.0
                    )
                elif ext == '.docx':
                    text = await asyncio.get_event_loop().run_in_executor(
                        _extraction_executor,
                        self._extract_from_docx,
                        temp_file_path
                    )
                elif ext in ['.pptx', '.ppt']:
                    text = await asyncio.get_event_loop().run_in_executor(
                        _extraction_executor,
                        self._extract_from_pptx,
                        temp_file_path
                    )
                elif ext in ['.txt', '.md']:
                    text = self._extract_from_text(temp_file_path)
                else:
                    logger.error(f"Unsupported file type: {ext}")
                    return None
                    
                logger.info(f"Extracted {len(text):,} characters from {ext}")
                
            except asyncio.TimeoutError:
                logger.error(f"Timeout extracting text from {file_path} ({file_size_mb:.2f} MB) - extraction took longer than 90 seconds")
                return None
            
            # Cleanup temp file
            if os.path.exists(temp_file_path):
                os.remove(temp_file_path)
            
            return text
            
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            # Cleanup on error
            if os.path.exists(temp_file_path):
                os.remove(temp_file_path)
            return None
    
    def _extract_from_pdf(self, file_path: str) -> str:
        """Extract text from PDF using PyMuPDF (faster) or PyPDF2 fallback"""
        try:
            if PYMUPDF_AVAILABLE:
                # PyMuPDF is 5-10x faster for large PDFs
                doc = fitz.open(file_path)
                page_count = len(doc)
                text = ""
                for page_num in range(page_count):
                    page = doc[page_num]
                    text += page.get_text() + "\n\n"
                doc.close()
                logger.info(f"Extracted PDF using PyMuPDF ({page_count} pages)")
                return text.strip()
            else:
                # Fallback to PyPDF2
                from PyPDF2 import PdfReader
                reader = PdfReader(file_path)
                page_count = len(reader.pages)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n\n"
                logger.info(f"Extracted PDF using PyPDF2 ({page_count} pages)")
                return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PDF: {e}")
            raise
    
    def _extract_from_docx(self, file_path: str) -> str:
        """Extract text from DOCX"""
        try:
            doc = Document(file_path)
            text = "\n\n".join([para.text for para in doc.paragraphs if para.text.strip()])
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting DOCX: {e}")
            raise
    
    def _extract_from_pptx(self, file_path: str) -> str:
        """Extract text from PPTX"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            return text.strip()
        except Exception as e:
            logger.error(f"Error extracting PPTX: {e}")
            raise
    
    def _extract_from_text(self, file_path: str) -> str:
        """Extract text from TXT/MD files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read().strip()
        except Exception as e:
            logger.error(f"Error extracting text file: {e}")
            raise

# Global document loader instance
document_loader = DocumentLoader()
